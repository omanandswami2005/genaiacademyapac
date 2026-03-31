"""Build the Gen AI Academy APAC hackathon presentation as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette ─────────────────────────────────────────────────────────
GOOGLE_BLUE   = RGBColor(0x42, 0x85, 0xF4)
GOOGLE_RED    = RGBColor(0xEA, 0x43, 0x35)
GOOGLE_YELLOW = RGBColor(0xFB, 0xBC, 0x05)
GOOGLE_GREEN  = RGBColor(0x34, 0xA8, 0x53)
DARK_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
CARD_BG       = RGBColor(0x16, 0x2A, 0x3E)   # card navy
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xC8, 0xD8, 0xE8)
ACCENT        = RGBColor(0x00, 0xC8, 0xFF)   # cyan


def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Google Sans" if False else "Calibri"
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=13, title=None, title_color=ACCENT,
                   bullet_color=WHITE, title_size=15):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Calibri"

    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = f"  •  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = bullet_color
        run.font.name = "Calibri"
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             accent_color=GOOGLE_BLUE):
    """Coloured card with title and bullet lines."""
    card = add_rect(slide, left, top, width, height, CARD_BG)

    # accent top border
    add_rect(slide, left, top, width, 0.04, accent_color)

    # title
    add_text_box(slide, title,
                 left + 0.15, top + 0.08, width - 0.3, 0.35,
                 font_size=13, bold=True, color=accent_color)

    # body
    y = top + 0.45
    for line in body_lines:
        add_text_box(slide, f"• {line}",
                     left + 0.2, y, width - 0.4, 0.32,
                     font_size=11, color=LIGHT_GREY)
        y += 0.30
    return card


# ═══════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank

# ───────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE SLIDE
# ───────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, DARK_BG)

# gradient-like accent bar on left
add_rect(s1, 0, 0, 0.35, 7.5, GOOGLE_BLUE)
add_rect(s1, 0.35, 0, 0.04, 7.5, ACCENT)

# Hackathon badge top-right
badge = add_rect(s1, 9.8, 0.25, 3.0, 0.5, GOOGLE_BLUE)
add_text_box(s1, "Gen AI Academy APAC Hackathon  |  Track 2",
             9.85, 0.28, 2.9, 0.44, font_size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

# Main title
add_text_box(s1, "🌍  Travel Guide Agent",
             0.7, 1.4, 9.5, 1.2,
             font_size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Cyan underline
add_rect(s1, 0.7, 2.65, 6.5, 0.06, ACCENT)

# Subtitle
add_text_box(s1,
             "AI-powered travel assistant using Google ADK + MCP\n"
             "Connected to real-world APIs via Model Context Protocol",
             0.7, 2.85, 9.5, 1.0,
             font_size=18, bold=False, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

# Tech pills row
pills = [
    ("Google ADK", GOOGLE_BLUE),
    ("Gemini 2.5 Flash", GOOGLE_GREEN),
    ("MCP", ACCENT),
    ("Cloud Run", GOOGLE_YELLOW),
    ("Vertex AI", GOOGLE_RED),
]
x = 0.7
for label, col in pills:
    w = len(label) * 0.095 + 0.4
    pill = add_rect(s1, x, 4.05, w, 0.38, col)
    add_text_box(s1, label, x + 0.05, 4.1, w - 0.1, 0.3,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += w + 0.18

# Live demo URL
add_text_box(s1,
             "🔗  Live:  https://travel-guide-agent-rnedu666qq-uc.a.run.app/dev-ui/",
             0.7, 4.65, 10.5, 0.4,
             font_size=12, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

# GitHub
add_text_box(s1,
             "⭐  GitHub:  github.com/omanandswami2005/genaiacademyapac",
             0.7, 5.1, 10.5, 0.4,
             font_size=12, bold=False, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

# Author
add_text_box(s1, "Omanand Swami  |  March 2026",
             0.7, 6.8, 6.0, 0.4,
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

# decorative globe emoji big
add_text_box(s1, "🌐", 10.8, 3.2, 2.0, 2.0,
             font_size=90, color=WHITE, align=PP_ALIGN.CENTER)


# ───────────────────────────────────────────────────────────────────────────
# SLIDE 2 — PROBLEM & SOLUTION
# ───────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, DARK_BG)

add_rect(s2, 0, 0, 13.33, 0.08, GOOGLE_BLUE)
add_text_box(s2, "The Problem & Our Solution",
             0.5, 0.2, 12.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_rect(s2, 0.5, 0.85, 4.0, 0.04, ACCENT)

# LEFT — Problem card
add_card(s2, 0.4, 1.1, 5.8, 5.5,
         "❌  The Problem",
         [
             "Travel info is scattered across 10+ sources",
             "Tourist guides are static & often outdated",
             "LLMs hallucinate travel facts (currencies,",
             "   visas, capitals, population figures)",
             "No standardised way for AI agents to call",
             "   external tools & APIs in real-time",
         ],
         accent_color=GOOGLE_RED)

# RIGHT — Solution card
add_card(s2, 6.8, 1.1, 5.8, 5.5,
         "✅  Our Solution",
         [
             "ADK agent backed by Gemini 2.5 Flash",
             "MCP server as single source of truth",
             "REST Countries API → authoritative data",
             "Wikipedia API → rich contextual summaries",
             "Region-aware travel advisory generator",
             "Deployed on Cloud Run — zero cold-start myth",
         ],
         accent_color=GOOGLE_GREEN)

# divider arrow
add_text_box(s2, "→", 6.2, 3.3, 0.6, 0.8,
             font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ───────────────────────────────────────────────────────────────────────────
# SLIDE 3 — ARCHITECTURE
# ───────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, DARK_BG)

add_rect(s3, 0, 0, 13.33, 0.08, GOOGLE_BLUE)
add_text_box(s3, "Architecture — How MCP Connects Everything",
             0.5, 0.2, 12.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_rect(s3, 0.5, 0.85, 5.5, 0.04, ACCENT)

# Arch flow boxes
boxes = [
    (0.4,  1.2, 2.5, 1.4, "👤  User",         "Ask: 'Tell me about Japan'",    GOOGLE_BLUE),
    (3.4,  1.2, 2.8, 1.4, "🤖  ADK Agent",    "Gemini 2.5 Flash\nVertex AI",  GOOGLE_GREEN),
    (6.7,  1.2, 2.8, 1.4, "🔧  MCP Server",   "3 Tools via stdio", ACCENT),
    (0.4,  3.5, 2.5, 1.4, "🌐  REST Countries","v3.1 API\nFacts & Stats",       GOOGLE_YELLOW),
    (3.4,  3.5, 2.8, 1.4, "📖  Wikipedia API", "REST summary\nendpoint",        GOOGLE_RED),
    (6.7,  3.5, 2.8, 1.4, "🏖️  Travel Tips",  "Region-aware\nadvisory gen",    GOOGLE_GREEN),
]

# Draw boxes
for bx, by, bw, bh, title, sub, col in boxes:
    add_rect(s3, bx, by, bw, bh, CARD_BG)
    add_rect(s3, bx, by, bw, 0.05, col)
    add_text_box(s3, title, bx+0.1, by+0.08, bw-0.2, 0.4,
                 font_size=13, bold=True, color=col)
    add_text_box(s3, sub, bx+0.1, by+0.52, bw-0.2, 0.7,
                 font_size=11, color=LIGHT_GREY)

# Arrows between flow boxes (top row)
for ax in [2.95, 6.25]:
    add_text_box(s3, "→", ax, 1.68, 0.45, 0.5,
                 font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Down arrows from MCP to APIs
for ax in [0.4+1.25-0.2, 3.4+1.4-0.2, 6.7+1.4-0.2]:
    add_text_box(s3, "↓", ax, 2.72, 0.45, 0.6,
                 font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Right side — MCP detail
add_rect(s3, 10.0, 1.1, 3.0, 5.5, CARD_BG)
add_rect(s3, 10.0, 1.1, 3.0, 0.05, ACCENT)
add_text_box(s3, "MCP Protocol Details",
             10.1, 1.18, 2.8, 0.4, font_size=12, bold=True, color=ACCENT)

mcp_details = [
    "Transport: stdio (subprocess)",
    "get_country_info()",
    "  → REST Countries v3.1",
    "  → name, capital, pop,",
    "    lang, currency, flag",
    "",
    "get_wikipedia_summary()",
    "  → /api/rest_v1/page/",
    "    summary/{topic}",
    "",
    "get_travel_advisory()",
    "  → Rule-based generator",
    "  → Region-specific tips",
]
y = 1.65
for line in mcp_details:
    col = ACCENT if "()" in line else LIGHT_GREY
    add_text_box(s3, line, 10.15, y, 2.7, 0.28,
                 font_size=10, color=col)
    y += 0.28


# ───────────────────────────────────────────────────────────────────────────
# SLIDE 4 — MCP DEEP DIVE
# ───────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, DARK_BG)

add_rect(s4, 0, 0, 13.33, 0.08, GOOGLE_GREEN)
add_text_box(s4, "Why MCP?  —  The Core Innovation",
             0.5, 0.2, 12.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_rect(s4, 0.5, 0.85, 4.5, 0.04, GOOGLE_GREEN)

# What is MCP block
add_rect(s4, 0.4, 1.1, 5.6, 2.8, CARD_BG)
add_rect(s4, 0.4, 1.1, 5.6, 0.05, GOOGLE_GREEN)
add_text_box(s4, "What is Model Context Protocol?",
             0.55, 1.18, 5.3, 0.4, font_size=13, bold=True, color=GOOGLE_GREEN)
paras = [
    "• Open standard by Anthropic, adopted by Google",
    "• Standardises how AI agents call external tools",
    "• Like USB-C for AI: one protocol, any tool",
    "• Agent stays decoupled from data sources",
    "• Works via JSON-RPC over stdio / HTTP / SSE",
]
y = 1.7
for p in paras:
    add_text_box(s4, p, 0.55, y, 5.2, 0.3, font_size=11.5, color=LIGHT_GREY)
    y += 0.35

# Without vs With MCP
add_card(s4, 0.4, 4.1, 5.6, 2.8,
         "Without MCP",
         [
             "Every tool needs custom integration code",
             "Agent code tangled with API logic",
             "Hard to swap or add tools",
             "No standard error handling or schemas",
         ], accent_color=GOOGLE_RED)

add_card(s4, 6.8, 4.1, 5.6, 2.8,
         "With MCP  ✅",
         [
             "Tools are plug-and-play modules",
             "Agent just calls tool by name + args",
             "Add new tools without touching agent",
             "Typed schemas, standard lifecycle",
         ], accent_color=GOOGLE_GREEN)

# Right — ADK + MCP code snippet
add_rect(s4, 6.8, 1.1, 6.1, 2.8, RGBColor(0x0A, 0x14, 0x1E))
add_text_box(s4, "ADK Agent Code  (agent.py)",
             6.95, 1.15, 5.8, 0.35, font_size=11, bold=True, color=ACCENT)
code = (
    "root_agent = LlmAgent(\n"
    "  model='gemini-2.5-flash',\n"
    "  name='travel_guide_agent',\n"
    "  tools=[\n"
    "    McpToolset(\n"
    "      connection_params=\n"
    "        StdioConnectionParams(\n"
    "          server_params=\n"
    "            StdioServerParameters(\n"
    "              command=sys.executable,\n"
    "              args=[MCP_SERVER_PATH]\n"
    "            )\n"
    "        )\n"
    "    )\n"
    "  ]\n"
    ")"
)
add_text_box(s4, code, 6.9, 1.55, 5.8, 2.3,
             font_size=9.5, color=GOOGLE_GREEN)


# ───────────────────────────────────────────────────────────────────────────
# SLIDE 5 — LIVE DEMO & RESULTS
# ───────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, DARK_BG)

add_rect(s5, 0, 0, 13.33, 0.08, GOOGLE_YELLOW)
add_text_box(s5, "Live Demo & Real Results",
             0.5, 0.2, 12.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_rect(s5, 0.5, 0.85, 3.5, 0.04, GOOGLE_YELLOW)

# Demo query
add_rect(s5, 0.4, 1.1, 12.5, 0.55, RGBColor(0x1A, 0x1A, 0x2E))
add_text_box(s5, "User asks:  'Tell me about Japan'",
             0.6, 1.17, 12.0, 0.4,
             font_size=14, bold=True, color=GOOGLE_YELLOW)

# Step flow
steps = [
    ("1", GOOGLE_BLUE,   "get_country_info('Japan')",
     "Capital: Tokyo  |  Pop: 123M  |  Currency: JPY  |  Language: Japanese  |  Region: Asia"),
    ("2", GOOGLE_GREEN,  "get_wikipedia_summary('Japan')",
     "'Japan is an island country in East Asia... Tokyo is the country's capital and largest city'"),
    ("3", GOOGLE_RED,    "get_travel_advisory('Japan')",
     "Located in Eastern Asia • Currency: Japanese Yen • Pop: 123M • 'Check visa requirements in advance'"),
    ("4", ACCENT,        "Agent combines & responds",
     "Rich, formatted travel guide with facts, context, and practical tips — all grounded in live data"),
]
y = 1.85
for num, col, title, result in steps:
    # number circle
    add_rect(s5, 0.4, y, 0.45, 0.45, col)
    add_text_box(s5, num, 0.4, y+0.02, 0.45, 0.4,
                 font_size=16, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text_box(s5, title, 0.95, y+0.03, 3.5, 0.38,
                 font_size=12, bold=True, color=col)
    add_text_box(s5, result, 4.55, y+0.03, 8.5, 0.38,
                 font_size=10.5, color=LIGHT_GREY)
    add_rect(s5, 0.4, y+0.5, 12.5, 0.01, CARD_BG)
    y += 0.6

# Metrics row
metrics = [
    ("3", "MCP Tools", GOOGLE_BLUE),
    ("2", "Live APIs", GOOGLE_GREEN),
    ("< 5s", "Response Time", GOOGLE_YELLOW),
    ("Cloud Run", "Serverless", GOOGLE_RED),
    ("Gemini 2.5", "Flash LLM", ACCENT),
]
x = 0.4
for val, label, col in metrics:
    w = 2.3
    add_rect(s5, x, 5.5, w, 1.1, CARD_BG)
    add_rect(s5, x, 5.5, w, 0.04, col)
    add_text_box(s5, val, x, 5.58, w, 0.5,
                 font_size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text_box(s5, label, x, 6.1, w, 0.35,
                 font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    x += w + 0.17

# Live URL badge
add_rect(s5, 3.5, 6.9, 6.5, 0.4, GOOGLE_BLUE)
add_text_box(s5,
             "🔗  https://travel-guide-agent-rnedu666qq-uc.a.run.app/dev-ui/",
             3.55, 6.93, 6.4, 0.35,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ───────────────────────────────────────────────────────────────────────────
# SLIDE 6 — TECH DEEP DIVE
# ───────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, DARK_BG)

add_rect(s6, 0, 0, 13.33, 0.08, ACCENT)
add_text_box(s6, "Technical Implementation",
             0.5, 0.2, 12.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_rect(s6, 0.5, 0.85, 3.5, 0.04, ACCENT)

tech_cards = [
    (0.4, 1.1, 3.8, 3.0, "ADK Agent Layer", [
        "LlmAgent with McpToolset",
        "Gemini 2.5 Flash via Vertex AI",
        "System instruction with tool usage guide",
        "FastAPI server via get_fast_api_app()",
        "SQLite session persistence (aiosqlite)",
    ], GOOGLE_BLUE),
    (4.5, 1.1, 3.8, 3.0, "MCP Server Layer", [
        "mcp.server.lowlevel.Server",
        "stdio transport (subprocess)",
        "3 tools with JSON schemas",
        "httpx async HTTP client",
        "Proper User-Agent for Wikipedia",
    ], ACCENT),
    (8.6, 1.1, 4.4, 3.0, "Infrastructure Layer", [
        "Cloud Run (us-central1) — serverless",
        "python:3.12-slim Docker image",
        "Non-root container user (security)",
        "1Gi memory, 300s timeout",
        "Artifact Registry + Cloud Build",
        "Vertex AI (GOOGLE_GENAI_USE_VERTEXAI=True)",
    ], GOOGLE_GREEN),
]

for cx, cy, cw, ch, title, items, col in tech_cards:
    add_card(s6, cx, cy, cw, ch, title, items, col)

# Bottom — key decisions
add_text_box(s6, "Key Engineering Decisions",
             0.5, 4.3, 12.0, 0.45,
             font_size=15, bold=True, color=ACCENT)

decisions = [
    ("stdio MCP transport", "Simpler than HTTP-SSE; agent and server co-locate in same container"),
    ("sys.executable as command", "Ensures MCP server uses the same Python/venv as the agent"),
    ("Synchronous root_agent", "Required for Cloud Run: async agents cause deployment issues in ADK"),
    ("Wikipedia User-Agent", "Wikimedia bot policy requires contact email to avoid 403 responses"),
]

y = 4.9
for title, desc in decisions:
    add_text_box(s6, f"› {title}:", 0.5, y, 3.2, 0.35,
                 font_size=11, bold=True, color=GOOGLE_YELLOW)
    add_text_box(s6, desc, 3.8, y, 9.0, 0.35,
                 font_size=11, color=LIGHT_GREY)
    y += 0.42


# ───────────────────────────────────────────────────────────────────────────
# SLIDE 7 — WHAT I LEARNED + FUTURE
# ───────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank_layout)
set_bg(s7, DARK_BG)

add_rect(s7, 0, 0, 13.33, 0.08, GOOGLE_GREEN)
add_text_box(s7, "Learnings & Future Roadmap",
             0.5, 0.2, 12.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_rect(s7, 0.5, 0.85, 4.0, 0.04, GOOGLE_GREEN)

# Learnings
add_card(s7, 0.4, 1.1, 5.8, 5.0,
         "🎓  What I Learned",
         [
             "MCP = 'USB-C for AI tools' — plug any",
             "  data source into any ADK agent",
             "",
             "ADK structures agent dev like software",
             "  engineering — clear separation of concerns",
             "",
             "Cloud Run + ADK = zero-ops deployment:",
             "  just gcloud run deploy --source .",
             "",
             "Wikimedia API requires proper User-Agent",
             "  (contact email) — API etiquette matters",
             "",
             "Stdio MCP is simpler than HTTP-SSE for",
             "  single-container deployments",
         ], accent_color=GOOGLE_GREEN)

# Future roadmap
add_card(s7, 6.8, 1.1, 5.8, 5.0,
         "🚀  Future Enhancements",
         [
             "Flight & hotel data via Amadeus MCP tool",
             "Google Maps Places API integration",
             "Currency conversion (live rates)",
             "Weather forecast tool",
             "",
             "Multi-language responses",
             "Voice interaction via Gemini Live API",
             "Chrome extension as UI layer",
             "",
             "Swap stdio → HTTP-SSE for multi-agent",
             "  orchestration (A2A protocol)",
             "Add evaluation with ADK EvalSets",
         ], accent_color=GOOGLE_BLUE)

# Bottom CTA
add_rect(s7, 2.5, 6.5, 8.0, 0.65, GOOGLE_BLUE)
add_text_box(s7,
             "🌐  Try it live:  https://travel-guide-agent-rnedu666qq-uc.a.run.app/dev-ui/"
             "     ⭐  github.com/omanandswami2005/genaiacademyapac",
             2.55, 6.55, 7.9, 0.55,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ───────────────────────────────────────────────────────────────────────────
# SAVE
# ───────────────────────────────────────────────────────────────────────────
out = "TravelGuideAgent_Presentation.pptx"
prs.save(out)
print(f"✅  Saved: {out}  ({prs.slides.__len__()} slides)")
